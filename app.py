# -*- coding: utf-8 -*-
import streamlit as st

# --- 1. タイトル等のセッション状態の初期化 ---
if "app_title" not in st.session_state:
    st.session_state.app_title = "汎用AIFAQチャットシステム"
if "current_welcome_msg" not in st.session_state:
    st.session_state.current_welcome_msg = (
        "はじめまして！「汎用AIFAQ お助けAI」です。😊\n\n"
        "お手元の資料や操作マニュアル（Excel, PDF, Word, CSV, PowerPoint等）を左側のメニューからアップロードしていただければ、即座にその内容を学習して専門のお助けAIに進化します！\n\n"
        "再現したい成果物の出力サンプルもお持ちの場合は、そちらを読み込ませることで具体的なデータ作成手順を詳しくお調べします。何でもお気軽にご質問ください！"
    )

# --- 【最優先ルール】Streamlitのページ構成設定は、他のあらゆるコマンドより先に最上部で実行します ---
st.set_page_config(page_title=st.session_state.app_title, layout="wide")

import pandas as pd
import google.generativeai as genai
from google.api_core import exceptions  # Rate Limit(429) エラーを確実に捕捉するため
from docx import Document
from PyPDF2 import PdfReader
from pptx import Presentation
import io
import configparser
import os
import signal
import re
import time  # リトライ待機（スリープ）処理のため

# --- 2. APIキーの設定 (APIKEY.ini または クラウドのSecretsからハイブリッド取得) ---
def load_api_key():
    config = configparser.ConfigParser()
    file_path = 'APIKEY.ini'
    if os.path.exists(file_path):
        try:
            config.read(file_path, encoding='utf-8-sig')
            return config.get('GEMINI', 'API_KEY')
        except:
            pass
            
    try:
        if "GEMINI" in st.secrets and "API_KEY" in st.secrets["GEMINI"]:
            return st.secrets["GEMINI"]["API_KEY"]
        elif "API_KEY" in st.secrets:
            return st.secrets["API_KEY"]
    except:
        return None
        
    return None

INI_KEY = load_api_key()
EMBEDDED_API_KEY = INI_KEY


# --- 3. 各ファイル抽出関数 (マニュアル内のテーブル構造も合わせて精密に抽出) ---
def extract_from_docx(file):
    file.seek(0)
    doc = Document(file)
    paragraphs_text = "\n".join([para.text for para in doc.paragraphs])
    
    tables_text = []
    for t_idx, table in enumerate(doc.tables):
        tables_text.append(f"\n\n[資料内テーブル #{t_idx+1}]")
        for row in table.rows:
            row_data = [cell.text.strip().replace('\n', ' ') for cell in row.cells]
            tables_text.append(" | ".join(row_data))
            
    return paragraphs_text + "\n" + "\n".join(tables_text)

def extract_from_pdf(file):
    file.seek(0)
    reader = PdfReader(file)
    return "\n".join([page.extract_text() for page in reader.pages if page.extract_text()])

def extract_from_pptx(file):
    file.seek(0)
    prs = Presentation(file)
    text_runs = []
    for s_idx, slide in enumerate(prs.slides):
        text_runs.append(f"\n--- スライド #{s_idx+1} ---")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                text_runs.append(shape.text)
            if shape.has_table:
                table = shape.table
                text_runs.append(f"\n[スライド内テーブル]")
                for row in table.rows:
                    row_data = [cell.text.strip().replace('\n', ' ') for cell in row.cells]
                    text_runs.append(" | ".join(row_data))
    return "\n".join(text_runs)

def extract_from_excel(file):
    file.seek(0)
    all_sheets = pd.read_excel(file, sheet_name=None)
    text_data = []
    for sheet_name, df in all_sheets.items():
        text_data.append(f"--- シート名 (参照テーブル): {sheet_name} ---\n{df.to_string(index=False)}")
    return "\n".join(text_data)

def extract_from_csv(file):
    file.seek(0)
    try:
        df = pd.read_csv(file)
        return df.to_string(index=False)
    except:
        try:
            file.seek(0)
            df = pd.read_csv(file, encoding="shift-jis")
            return df.to_string(index=False)
        except Exception as e:
            return f"[CSV読込エラー] {str(e)}"

def extract_from_text(file):
    file.seek(0)
    try:
        return file.read().decode("utf-8")
    except:
        try:
            file.seek(0)
            return file.read().decode("shift-jis")
        except Exception as e:
            return f"[テキスト読込エラー] {str(e)}"


# --- 404エラーを回避しつつ、利用可能なモデル名を安全に取得する関数 ---
def get_safe_model_name(api_key):
    try:
        genai.configure(api_key=api_key)
        available_models = [m.name for m in genai.list_models() if 'generateContent' in m.supported_generation_methods]
        
        has_flash = any('gemini-1.5-flash' in m for m in available_models)
        target_raw = 'gemini-1.5-flash' if has_flash else (available_models[0] if available_models else 'gemini-1.5-flash')
        
        safe_name = target_raw.replace('models/', '')
        return safe_name
    except:
        return 'gemini-1.5-flash'

# ファイル名からメインのシステム・サービス名を抽出するクレンジング関数
def clean_service_name(filename):
    base = os.path.splitext(filename)[0]
    patterns = [
        r"(操作)?マニュアル", r"取扱説明書", r"手順書", r"仕様書", r"概要書", r"説明書",
        r"【.*】", r"\[.*\]", r"（.*）", r"\(.*\)",
        r"[vV]er\.?\d+(\.\d+)*", r"\d{8}", r"\d{6}",
        r"[-_]"
    ]
    cleaned = base
    for pat in patterns:
        cleaned = re.sub(pat, " ", cleaned)
    cleaned = cleaned.strip()
    return cleaned if cleaned else base

# マニュアルテキストを解析して、メインの製品・システム名をGeminiからスマートに特定する関数
def extract_service_name_via_ai(text, default_name, api_key):
    try:
        genai.configure(api_key=api_key)
        target_model = get_safe_model_name(api_key)
        model = genai.GenerativeModel(target_model)
        
        prompt = f"""
以下は、ユーザーからアップロードされたマニュアルまたは資料テキストの冒頭部分です。
この資料が「何のツール」「何のサービス」または「どのシステム」について説明しているものか、最もメインとなる固有名称を日本語で1つだけ見つけ出してください。
余計な説明、前置き、記号、拡張子などは絶対に含めず、純粋な名称のみを返してください。
最大でも20文字以内とします。特定が難しい場合は「{default_name}」を返してください。

【資料テキストの一部】
{text[:1500]}
"""
        response = model.generate_content(prompt)
        res_text = response.text.strip()
        res_text = re.sub(r"[`'\"]", "", res_text)
        res_text = res_text.split("\n")[0].strip()
        return res_text if res_text else default_name
    except:
        return default_name


# --- 4. AI回答生成ロジック (自動リトライ・履歴ウィンドウ削減版 / 汎用化プロンプト) ---
def get_ai_roleplay_response(messages, persona, product_docs, format_docs, api_key):
    target_model = get_safe_model_name(api_key)
    recent_messages = [messages[0]] + messages[-5:] if len(messages) > 6 else messages

    for attempt in range(5):
        try:
            genai.configure(api_key=api_key)
            model = genai.GenerativeModel(target_model)
            
            combined_docs = "\n\n".join(product_docs) if product_docs else "追加のマニュアル等のアップロードは現在ありません。一般的な知識に基づいてユーザーの課題解決を支援してください。"
            combined_formats = "\n\n".join(format_docs) if format_docs else "出力フォーマットサンプルの指定は現在ありません。"
            
            history_text = ""
            for m in recent_messages:
                role_label = "AIアシスタント(あなた)" if m["role"] == "assistant" else "ユーザー"
                history_text += f"{role_label}: {m['content']}\n"

            system_prompt = f"""
あなたはユーザーから提供されたマニュアルや資料に基づいて、操作手順や記載内容を正確に説明する熟練のAIFAQアシスタントです。
操作方法、機能説明、データ構造、エラー解決等に関する質問に対して、丁寧かつ極めて分かりやすく回答してください。

【対象マニュアル・参照資料情報】
{persona['description']}

【アップロードされた各種マニュアル・参考ドキュメント（最優先参照情報）】
{combined_docs}

【アップロードされた出力フォーマットサンプル】
{combined_formats}
※この出力フォーマットサンプルが提示されている場合は、ユーザーが「これと同じデータ形式やレイアウトを出力したい」と希望しています。
現在アップロードされている各種マニュアル・参考ドキュメントを参照し、このフォーマットを出力するにはどのような操作, 設定, データの選択や加工手順を行えばいいのかを、手順を追って具体的に説明してください。

【テーブル（表）の参照と特定に関する絶対ルール】
1. **「参照テーブルの明記」**: アップロードされた資料には、操作手順や設定値が「表（テーブル形式）」で整理されている箇所が多数あります。ユーザーの質問に答える際、または手順を解説する際には、マニュアル内の **どのテーブル（例：[資料内テーブル #1]、[シート名/テーブル名]、列項目名など）を参照してその判断や数値・手順を導き出したのか** を、回答内で必ず具体的に言及・特定してください。
2. データ項目やコード値の説明時には、「〇〇マニュアルの、表『××』に記載のある通り…」のように, ユーザーが自身でマニュアルを手繰って検証・確認できるように参照元テーブルを紐づけてください。

【預かり資産トータルクエリーサービスに関する絶対判定ルール】
もしアップロードされた資料の内容や質問の文脈が「預かり資産トータルクエリーサービス」に関連する場合、ユーザーのやりたいデータ抽出や操作要望に対して、以下の思考プロセスを厳格に適用して回答を構成してください。
1. **「標準クエリ（約200種類）の確認」**: まず第一に、ユーザーのやりたい要件にそのまま合致する既存の標準クエリがすでに提供されているかをマニュアル内のクエリ一覧テーブル等から判断して案内してください。
2. **「既存クエリの修正・加工方法」**: もし完全に合致する既存 of 標準クエリがそのままでは見つからない場合、どの標準クエリをベース（ひな形）に選択し、それをどのように修正（項目追加、結合、フィルター条件の編集など）すれば目的の結果が得られるかを、具体的かつ分かりやすい手順として説明してください。

【回答の絶対ルール】
1. ユーザーの質問に対し、アップロードされたマニュアルの情報を最も信頼できる「絶対の基準（最優先情報）」として参照し、正確に回答を構成してください。
2. アップロードされた情報だけで判断がつかない不確実な事項やマニュアルに記載がない操作については、知ったかぶりをせず、「マニュアル等に記載がありませんでした」と明示したうえで、一般的な推奨方法を補足するか、専門の窓口や管理者への確認を案内してください。
3. 専門用語が使われている場合でも、操作担当者がスムーズに迷わず作業を進められるよう、ステップ・バイ・ステップの具体的な手順や丁寧な表現で回答してください。
4. AIとしてのメタな発言（例：「以上がマニュアルに基づく回答です」など）は含めず、ユーザーへの親切な回答テキストのみを親みやすいLINE風の対話形式で出力してください。

【これまでの会話履歴】
{history_text}

Above rules and history will be used to generate the next response.
"""
            response = model.generate_content(system_prompt)
            return response.text

        except (exceptions.ResourceExhausted, Exception) as e:
            error_msg = str(e)
            if "429" in error_msg or isinstance(e, exceptions.ResourceExhausted):
                wait_time = (attempt + 1) * 10
                time.sleep(wait_time)
                continue
            return f"【システムエラー】詳細: {error_msg}"
            
    return "【混雑エラー】現在AIへのリクエストが連続しています。無料枠の制限を超過したため、1分ほど待ってから再度送信してください。"


# --- Rerun処理 of 安全な抽象化 ---
def safe_rerun():
    try:
        st.rerun()
    except AttributeError:
        try:
            st.experimental_rerun()
        except AttributeError:
            pass


# --- 5. アプリケーションのメインロジック ---
def main_app():
    if st.session_state.get("app_terminated", False):
        st.warning("🛑 システムは終了しました。再度ご利用になる場合は、ブラウザをリロード（再読み込み）してください。")
        st.stop()

    # ユニバーサルデザインに準拠した大きめの常時特大フォント（22px基準）
    base_font_px = "22px"
    bubble_font_px = "22px"

    # --- 🎨 温かみ・リッチ感・文字や各種メニューボタンの「超高コントラスト視認性」をCSSで完全追求 ---
    st.markdown(f"""<style>
/* 1. 標準サイドバー領域を強制的に100%排除し、コンテナ崩れを物理解消 */
[data-testid="stSidebar"] {{
display: none !important;
}}
[data-testid="stSidebarCollapsedControl"] {{
display: none !important;
}}

/* 画面全体の余白を最大化 & 見切れ対策のパディング追加 */
.block-container {{
padding-top: 3.5rem !important;
padding-bottom: 6rem !important;
padding-left: 2.5rem !important;
padding-right: 2.5rem !important;
max-width: 100% !important;
}}

/* アプリケーション全体背景 (温かみのあるニュアンスベージュ・オフホワイト) */
.stApp {{
background-color: #f6f5f0 !important;
font-family: "BIZ UDゴシック", "BIZ UDPゴシック", "Helvetica Neue", Arial, sans-serif !important;
font-size: {base_font_px} !important;
line-height: 1.6 !important;
color: #1a2a1d !important; /* 視認性を極限まで高めた濃いインクグリーン */
}}

/* Streamlitの標準 st.container の枠線を「お手本カード」に完全上書き (丸みと柔らかな陰影) */
div[data-testid="stVerticalBlockBorderWrapper"] {{
background-color: #ffffff !important;
border: 1.5px solid #d0d8cd !important; /* 温かみのあるオリーブグレーの太めボーダー */
border-radius: 14px !important;
box-shadow: 0 6px 12px -2px rgba(27,45,32,0.06) !important; /* 緑に調和したリッチな陰影 */
padding: 24px !important;
}}

/* チャットボックス＆吹き出しレイアウト */
.chat-container {{
display: flex;
flex-direction: column;
gap: 20px;
width: 100%;
margin-bottom: 25px;
}}
.chat-row-user {{
display: flex;
justify-content: flex-end;
width: 100%;
}}
.chat-row-assistant {{
display: flex;
justify-content: flex-start;
width: 100%;
}}
.chat-bubble-user {{
background-color: #f0f3ee; /* 薄いオリーブグレー */
color: #132817;
padding: 22px 28px;
border-radius: 20px 20px 0px 20px;
max-width: 82%;
border: 1.5px solid #c8d3c5;
font-size: {bubble_font_px} !important;
line-height: 1.6;
box-shadow: 0 3px 6px rgba(0,0,0,0.02);
}}
.chat-bubble-assistant {{
background-color: #ffffff;
color: #132817;
padding: 22px 28px;
border-radius: 20px 20px 20px 0px;
max-width: 82%;
border: 1.5px solid #c8d3c5;
font-size: {bubble_font_px} !important;
line-height: 1.6;
box-shadow: 0 3px 6px rgba(0,0,0,0.02);
}}

/* アバター用サークル（2倍サイズ化：68px） */
.avatar-circle-ai {{
width: 68px !important;
height: 68px !important;
border-radius: 50% !important;
background-color: #2d4a34 !important; /* より深みのあるリッチな常盤色 */
color: white !important;
display: flex !important;
align-items: center !important;
justify-content: center !important;
font-weight: bold !important;
font-size: 24px !important;
box-shadow: 0 4px 8px rgba(0,0,0,0.08) !important;
}}

/* チャットヘッダー（アバターと表記の一体化） */
.chat-header {{
display: flex;
align-items: center;
gap: 14px;
font-size: 20px;
font-weight: bold;
color: #132817;
}}

/* 送信ボタン等 (温かみのあるテラコッタオレンジ：#d36a43) */
div.stButton > button[type="submit"], div.stForm button {{
background-color: #d36a43 !important; /* コントラストを高めたテラコッタ */
color: white !important;
border-radius: 12px !important;
border: none !important;
padding: 0.9rem 2rem !important;
font-size: 20px !important;
font-weight: 800 !important;
width: 100% !important;
transition: all 0.2s !important;
box-shadow: 0 4px 10px rgba(211,106,67,0.25) !important;
}}
div.stButton > button[type="submit"]:hover, div.stForm button:hover {{
background-color: #ba572a !important;
box-shadow: 0 6px 14px rgba(211,106,67,0.35) !important;
}}

/* 【視認性の徹底強化】左パネルの見出し・メニューの文字色をハッキリした深緑に変更 */
.left-panel-title {{
font-size: 21px !important;
font-weight: 800 !important;
color: #1a331e !important;
margin-bottom: 12px !important;
display: flex;
align-items: center;
gap: 8px;
}}

/* ドラッグ＆ドロップ領域 (st.file_uploader) の高コントラスト枠線・視認性確保 */
div[data-testid="stFileUploader"] {{
border: 2px dashed #2d4a34 !important; /* 濃いオリーブ破線 */
background-color: #fafbfc !important;
border-radius: 12px !important;
padding: 15px !important;
}}
div[data-testid="stFileUploader"] section {{
color: #1a331e !important; /* 濃い文字 */
font-weight: bold !important;
}}
div[data-testid="stFileUploader"] button {{
background-color: #2d4a34 !important;
color: white !important;
font-weight: bold !important;
border-radius: 8px !important;
border: none !important;
box-shadow: 0 2px 4px rgba(0,0,0,0.1) !important;
}}

/* 入力エリアのフォントサイズアップ */
textarea, input {{
font-size: 20px !important;
color: #1a331e !important;
font-weight: 600 !important;
}}
</style>""", unsafe_allow_html=True)

    # --- 🌟 リッチモスグリーン（#2d4a34）超巨大ヘッダー (上部見切れを防ぐため十分な安全マージンを設定) ---
    st.markdown(f"""
    <div style="background-color: #2d4a34; padding: 25px 35px; border-radius: 14px; display: flex; align-items: center; gap: 20px; margin-top: 20px; margin-bottom: 25px; box-shadow: 0 6px 15px rgba(45,74,52,0.15); border: 1px solid #4a6c53;">
        <div style="background-color: #f6f5f0; border-radius: 50%; width: 72px; height: 72px; display: flex; align-items: center; justify-content: center; font-size: 38px; line-height: 1; box-shadow: inset 0 2px 4px rgba(0,0,0,0.05);">💡</div>
        <div>
            <div style="color: white; font-size: 34px; font-weight: 800; line-height: 1.1; letter-spacing: 0.5px;">{st.session_state.app_title}</div>
            <div style="color: #c9ded0; font-size: 16px; margin-top: 5px; font-weight: bold; letter-spacing: 0.8px;">公式マニュアル・お助けAIチャット</div>
        </div>
    </div>
    """, unsafe_allow_html=True)

    st.write("") # スペーサー

    # 左右スプリットレイアウト (左3: 右7) の完全メイン配置
    left_col, right_col = st.columns([3, 7])

    # --- 左側パネル（API設定、マニュアルアップローダー、出力サンプル） ---
    with left_col:
        # システム終了ボタンを目立ちやすく、視認性の高いテラコッタボタンに変更して配置
        if st.button("🛑 システムを終了する", key="btn_sys_term", use_container_width=True):
            st.session_state.app_terminated = True
            st.rerun()
            
        st.markdown("<div style='margin-top: 20px;'></div>", unsafe_allow_html=True)

        # 1. AIキー設定カード
        st.markdown('<div class="left-panel-title">🔑 AIキー設定（お持ちの場合のみ）</div>', unsafe_allow_html=True)
        with st.container(border=True):
            custom_api_key = st.text_input(
                "Gemini APIキーを入力",
                type="password",
                placeholder="入力すると優先して使われます",
                label_visibility="collapsed"
            )
            st.markdown("<div style='color: #1a331e; font-size: 15px; font-weight: bold; margin-top: 6px;'>※空欄の場合は、システム既定のキーで自動的に動くので安心してください。</div>", unsafe_allow_html=True)
            
            ACTIVE_API_KEY = custom_api_key if custom_api_key else EMBEDDED_API_KEY
            if ACTIVE_API_KEY:
                st.success("✔️ AIキーが有効に作動しています")
            else:
                st.warning("⚠️ APIキーを設定してください")

        st.markdown("<div style='margin-top: 20px;'></div>", unsafe_allow_html=True)

        # 2. マニュアル資料読込カード
        st.markdown('<div class="left-panel-title">📁 手元の資料をAIに読み込ませる</div>', unsafe_allow_html=True)
        with st.container(border=True):
            uploaded_files = st.file_uploader(
                "資料 (Excel, PDF, Word, CSV, PPT)",
                type=["docx", "pdf", "pptx", "xlsx", "xls", "csv"],
                accept_multiple_files=True,
                label_visibility="collapsed"
            )
            
            all_extra_text = []
            file_names = []
            if uploaded_files:
                for f in uploaded_files:
                    try:
                        if f.name.endswith(".docx"): content = extract_from_docx(f)
                        elif f.name.endswith(".pdf"): content = extract_from_pdf(f)
                        elif f.name.endswith(".pptx"): content = extract_from_pptx(f)
                        elif f.name.endswith((".xlsx", ".xls")): content = extract_from_excel(f)
                        elif f.name.endswith(".csv"): content = extract_from_csv(f)
                        else: content = ""
                        
                        if content:
                            all_extra_text.append(f"--- ファイル名: {f.name} ---\n{content}")
                            file_names.append(f.name)
                            st.success(f"✔️ {f.name} を学習完了")
                    except Exception as e:
                        st.error(f"❌ {f.name} 読み込み失敗: {str(e)}")

        # マニュアルに基づくダイナミックなタイトル＆ペルソナ変更
        if file_names:
            if "last_processed_files" not in st.session_state or st.session_state.last_processed_files != file_names:
                default_service_name = clean_service_name(file_names[0])
                if ACTIVE_API_KEY and all_extra_text:
                    joined_samples = "\n".join(all_extra_text)
                    detected_name = extract_service_name_via_ai(joined_samples, default_service_name, ACTIVE_API_KEY)
                    st.session_state.app_title = detected_name
                    st.session_state.current_welcome_msg = (
                        f"はじめまして！「{detected_name} お助けAI」です。😊\n\n"
                        f"アップロードされた資料【{', '.join(file_names)}】をすべて記憶して完全学習しました！\n"
                        f"具体的な操作手順、設定、活用方法など、何でも分かりやすくお答えします！"
                    )
                else:
                    st.session_state.app_title = default_service_name
                st.session_state.last_processed_files = file_names
                st.rerun()

            current_persona = {
                "description": f"提供されたマニュアル「{', '.join(file_names)}」（対象システム/ツール: {st.session_state.app_title}）に精通した、専属の優秀なAIFAQ操作説明アシスタントです。"
            }
        else:
            st.session_state.app_title = "汎用AIFAQシステム"
            current_persona = {
                "description": "現在は特定のマニュアルは読み込まれていません。アップロードされる多種多様なマニュアルに沿って、操作方法や記載内容に回答する汎用お助けAIFAQアシスタントです。"
            }

        st.markdown("<div style='margin-top: 20px;'></div>", unsafe_allow_html=True)

        # 3. 出力フォーマットサンプル読込カード
        st.markdown('<div class="left-panel-title">📋 出力サンプルの読込</div>', unsafe_allow_html=True)
        with st.container(border=True):
            if "format_samples" not in st.session_state:
                st.session_state.format_samples = []
            if "format_file_names" not in st.session_state:
                st.session_state.format_file_names = []

            uploaded_format = st.file_uploader(
                "出力サンプル (Word, PDF, Excel, CSV, PPT, TXT)",
                type=["docx", "pdf", "pptx", "xlsx", "xls", "csv", "txt"],
                key="format_uploader"
            )

            if uploaded_format:
                f_name = uploaded_format.name
                if f_name not in st.session_state.format_file_names:
                    try:
                        if f_name.endswith(".docx"): content = extract_from_docx(uploaded_format)
                        elif f_name.endswith(".pdf"): content = extract_from_pdf(uploaded_format)
                        elif f_name.endswith(".pptx"): content = extract_from_pptx(uploaded_format)
                        elif f_name.endswith((".xlsx", ".xls")): content = extract_from_excel(uploaded_format)
                        elif f_name.endswith(".csv"): content = extract_from_csv(uploaded_format)
                        elif f_name.endswith(".txt"): content = extract_from_text(uploaded_format)
                        else: content = ""
                        
                        if content:
                            st.session_state.format_samples.append(f"--- サンプルファイル名: {f_name} ---\n{content}")
                            st.session_state.format_file_names.append(f_name)
                            st.success(f"✔️ サンプル「{f_name}」学習完了")
                    except Exception as e:
                        st.error(f"❌ {f_name} の読込失敗: {str(e)}")

            if st.session_state.format_file_names:
                st.markdown("<div style='color:#1a331e; font-weight:bold; margin-top:10px;'>📌 ロード中サンプル:</div>", unsafe_allow_html=True)
                for name in st.session_state.format_file_names:
                    st.write(f"・ {name}")
                if st.button("🗑️ サンプルをクリア", use_container_width=True):
                    st.session_state.format_samples = []
                    st.session_state.format_file_names = []
                    st.success("クリアしました。")
                    safe_rerun()

    # --- 右側パネル（ステータス、チャットスレッド、入力欄） ---
    with right_col:
        # AIステータスバー（アバターやフォントサイズを拡大）
        status_col1, status_col2 = st.columns([8, 2])
        with status_col1:
            st.markdown(f"""
            <div style="background-color: white; border: 1.5px solid #d0d8cd; border-radius: 12px; padding: 18px 24px; display: flex; align-items: center; gap: 16px; box-shadow: 0 4px 8px rgba(0,0,0,0.02);">
                <div class="avatar-circle-ai" style="width: 58px !important; height: 58px !important; font-size: 20px !important;">AI</div>
                <div>
                    <div style="font-weight: bold; font-size: 22px; color: #132817;">{st.session_state.app_title} サポートAI</div>
                    <div style="color: #3b5e43; font-size: 15px; font-weight: bold; margin-top: 3px;">● いつでも質問に答えますよ</div>
                </div>
            </div>
            """, unsafe_allow_html=True)
        with status_col2:
            st.markdown("""<div style="height:100%; display:flex; align-items:center; justify-content:center;">
<button style="border:1.5px solid #d0d8cd; background:white; border-radius:12px; padding:16px; width:100%; font-weight:bold; cursor:pointer; font-size:16px; color:#a0aec0; height: 100%;">🔇 音声: 停止中</button>
</div>""", unsafe_allow_html=True)

        st.write("")

        # 履歴セッションの初期化
        if "messages" not in st.session_state:
            st.session_state.messages = [{
                "role": "assistant",
                "content": st.session_state.current_welcome_msg
            }]

        # チャットログのレンダリング
        chat_placeholder = st.container()
        with chat_placeholder:
            st.markdown('<div class="chat-container">', unsafe_allow_html=True)
            for m in st.session_state.messages:
                row_class = "chat-row-assistant" if m["role"] == "assistant" else "chat-row-user"
                bubble_class = "chat-bubble-assistant" if m["role"] == "assistant" else "chat-bubble-user"
                
                if m["role"] == "assistant":
                    header_html = f'<div class="chat-header"><div class="avatar-circle-ai">AI</div><span style="font-size: 20px;">サポートAI</span></div>'
                else:
                    header_html = f'<div class="chat-header"><span style="font-size:36px; line-height: 1;">💼</span><span style="font-size: 20px;">あなた（入力者）</span></div>'

                st.markdown(f"""
                <div class="{row_class}">
                    <div class="{bubble_class}">
                        {header_html}
                        <div style="margin-top: 12px;">{m["content"]}</div>
                    </div>
                </div>
                """, unsafe_allow_html=True)
            st.markdown('</div>', unsafe_allow_html=True)

        # チャット入力処理（常に最下段にフローティング固定される st.chat_input を採用）
        if prompt := st.chat_input("ここに知りたいことを入力してください（例：AWSエラーが出る など）"):
            st.session_state.messages.append({"role": "user", "content": prompt})
            
            with chat_placeholder:
                st.markdown(f"""
                <div class="chat-row-user">
                    <div class="chat-bubble-user">
                        <div class="chat-header"><span style="font-size:36px; line-height: 1;">💼</span><span style="font-size: 20px;">あなた（入力者）</span></div>
                        <div style="margin-top: 12px;">{prompt}</div>
                    </div>
                </div>
                """, unsafe_allow_html=True)
                
            with st.chat_message("assistant", avatar="🤖"):
                with st.spinner("情報をマニュアルから参照して回答を作成中..."):
                    res = get_ai_roleplay_response(
                        st.session_state.messages,
                        current_persona,
                        all_extra_text,
                        st.session_state.format_samples,
                        ACTIVE_API_KEY
                    )
                    st.markdown(res)
                
            st.session_state.messages.append({"role": "assistant", "content": res})
            st.rerun()

# アプリを実行
main_app()
